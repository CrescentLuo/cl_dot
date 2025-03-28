import sys
import io
import os
import re
import json
import logging
import argparse
import shutil
import zipfile
import olefile
import binwalk
import pandas as pd
import altair as alt
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm


def process_text(text):
    try:
        with open(Path(__file__).parent / "dify_api_token.json") as f:
            token_data = json.load(f)
            api_token = token_data.get("api_token")
            if not api_token:
                logging.error("API token not found in token file")
                return None
    except Exception as e:
        logging.error(f"Failed to load API token: {e}")
        return None

    api_url = "http://10.10.1.29/v1/chat-messages"
    headers = {
        "Authorization": f"Bearer {api_token}",
        "Content-Type": "application/json",
    }

    payload = {
        "inputs": {},
        "query": text,
        "response_mode": "blocking",
        "user": "zheng.luo",
    }

    try:
        response = requests.post(api_url, headers=headers, json=payload)
        response.raise_for_status()
        response_data = response.json()

        # Extract information from the answer
        extracted_info = {}
        if "answer" in response_data:
            answer_lines = re.split(r"<br/>|\n", response_data["answer"])
            for line in answer_lines:
                line = line.strip()
                if ":" in line:
                    key, value = line.split(":", 1)
                    extracted_info[key.strip()] = value.strip()
            response_data["extracted_info"] = extracted_info
        return response_data
    except Exception as e:
        logging.error(f"API request failed: {e}")
        return None


def extract_portion(input_file, output_file, block_size, count, skip):
    start_pos = skip * block_size
    bytes_to_read = block_size * count

    with open(input_file, "rb") as f_in:
        f_in.seek(start_pos)
        data = f_in.read(bytes_to_read)

    with open(output_file, "wb") as f_out:
        f_out.write(data)


def extract_descriptions(file_path):
    descriptions = []
    offsets = []
    zip_archive = []

    # Collect all results
    for module in binwalk.scan(file_path, signature=True, quiet=True):
        for result in module.results:
            descriptions.append(result.description)
            offsets.append(result.offset)

    # Process each description
    output_file_set = set()
    for i in range(len(descriptions)):
        current_offset = offsets[i]
        next_offset = offsets[i + 1] if i + 1 < len(offsets) else "EOF"

        if "name: " in descriptions[i]:
            output_file = descriptions[i].split("name: ")[1].strip()
            if output_file not in output_file_set:
                output_file_set.add(output_file)
                zip_file = f"{file_path}_{i}.zip"
                zip_archive.append(zip_file)
                extract_portion(
                    file_path,
                    zip_file,
                    block_size=1,
                    count=next_offset - current_offset,
                    skip=current_offset,
                )
    return zip_archive


def binwalk_scan(file_in, quiet=True):
    try:
        zip_archive = extract_descriptions(file_in)
        destination_path = Path(file_in).parent / f"{Path(file_in).stem}"
        destination_path.mkdir(exist_ok=True)

        if zip_archive:
            for zip_file in tqdm(zip_archive, desc="Processing zip files"):
                try:
                    binwalk.scan(
                        str(zip_file),
                        signature=True,
                        extract=True,
                        rm=True,
                        quiet=quiet,
                        directory=str(destination_path),
                    )
                    Path(zip_file).unlink()

                    # Move extracted files from subfolder to parent folder
                    extracted_subfolder = (
                        destination_path / f"_{Path(zip_file).name}.extracted"
                    )
                    source_dir = str(extracted_subfolder)
                    destination_dir = str(destination_path)

                    # Use shutil instead of os.system
                    if extracted_subfolder.exists():
                        for item in extracted_subfolder.glob("*"):
                            if item.is_file():
                                shutil.copy2(item, destination_path)
                            elif item.is_dir():
                                shutil.copytree(
                                    item,
                                    destination_path / item.name,
                                    dirs_exist_ok=True,
                                )
                        shutil.rmtree(extracted_subfolder)

                except Exception as e:
                    logging.error(f"Error processing {zip_file}: {e}")
    except Exception as e:
        logging.error(f"Error in binwalk scan: {e}")


def extract_prism_files(filepath, single_slide=None):
    prs = Presentation(filepath)
    input_file_directory = Path(filepath).parent
    extracted_files = []

    slides = prs.slides
    if not single_slide:
        slides = tqdm(slides, desc="Processing slides", total=len(prs.slides))

    for slide in slides:
        slide_number = prs.slides.index(slide) + 1
        if single_slide and slide_number != single_slide:
            continue

        for shape in slide.shapes:
            if shape.shape_type in [
                MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
                MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            ]:
                shape_id = shape.shape_id
                ole = shape.ole_format
                blob = ole.blob
                bytes_io = io.BytesIO(blob)

                if zipfile.is_zipfile(bytes_io):
                    with olefile.OleFileIO(bytes_io) as ole_stream:
                        stream_name = "CONTENTS"
                        if ole_stream.exists(stream_name):
                            blob = ole_stream.openstream(stream_name).read()

                filename = (
                    input_file_directory / f"slide_{slide_number}_id_{shape_id}.bin"
                )
                with open(filename, "wb") as f:
                    f.write(blob)
                binwalk_scan(str(filename), quiet=True)
                extracted_files.append(str(filename))

    return extracted_files


class ExtractedBin:
    def __init__(self, file_path):
        self.file_path = Path(file_path)
        self.sheets = []
        self.tables = []
        self.extract_name()

    def extract_name(self):
        match = re.search(r"slide_(\d+)_id_(\d+)\.bin", str(self.file_path.name))
        if match:
            self.slide_number = match.group(1)
            self.shape_id = match.group(2)
        else:
            self.slide_number = "unknown"
            self.shape_id = "unknown"
        self.file_path = self.file_path.parent / f"{self.file_path.stem}"

    def extract_sheets(self):
        data_folder = Path(self.file_path) / "data"
        sheets_folder = data_folder / "sheets"

        if not sheets_folder.exists():
            return

        for sub_dir in sheets_folder.iterdir():
            sheet_json = sub_dir / "sheet.json"
            if sheet_json.exists():
                with open(sheet_json) as json_file:
                    self.sheets.append(json.load(json_file))

    def extact_table_dataSets(self, uid):
        sets_folder = self.file_path / "data" / "sets"
        sets_json_path = sets_folder / f"{uid}.json"

        with open(sets_json_path) as json_file:
            sets_meta = json.load(json_file)
            return {
                uid: {
                    "fenID": sets_meta["fenID"],
                    "setName": (
                        sets_meta["title"]["string"]
                        if isinstance(sets_meta["title"], dict)
                        else sets_meta["title"]
                    ),
                }
            }

    def extract_table(self):
        if not self.sheets:
            self.extract_sheets()

        for sheet_meta in self.sheets:
            table_uid = sheet_meta["table"]["uid"]
            dataSets = sheet_meta["table"]["dataSets"]
            repCnt = sheet_meta["table"]["replicatesCount"]

            header_row = []
            for set_uid in dataSets:
                set_meta = self.extact_table_dataSets(set_uid)
                for i in range(repCnt):
                    header_row.append(set_meta[set_uid]["setName"] + f"_rep{i}")

            table_file_path = (
                self.file_path / "data" / "tables" / table_uid / "data.csv"
            )

            csv_table = pd.read_csv(
                table_file_path,
                header=None,
                index_col=0,
                names=["LNP"] + header_row,
                usecols=range(len(header_row) + 1),
            )

            csv_table = csv_table.reset_index().melt(
                id_vars="LNP", var_name="Tissue", value_name="value"
            )

            csv_table[["Tissue", "rep"]] = csv_table["Tissue"].str.split(
                pat="_", expand=True
            )

            self.tables.append(csv_table)

    def save_tables(self):
        for idx, _ in enumerate(self.sheets):
            output_path = f"{self.file_path.parent}/Page_{self.slide_number}_Shape_{self.shape_id}_{idx}"
            self.tables[idx].to_csv(f"{output_path}.csv")
            self.plot_table(self.tables[idx]).save(f"{output_path}.svg")

    def plot_table(self, table):
        bar_chart = (
            alt.Chart(table)
            .mark_bar()
            .encode(
                x="LNP:O",
                xOffset="Tissue:N",
                y=alt.Y("mean(value):Q", scale=alt.Scale(type="log")).title(
                    "Mean value"
                ),
                color="Tissue:N",
            )
        )

        point_chart = (
            alt.Chart(table)
            .mark_circle()
            .encode(
                x="LNP:O",
                xOffset="Tissue:N",
                y=alt.Y("value:Q", scale=alt.Scale(type="log")),
                color=alt.value("black"),
            )
        )

        return alt.layer(bar_chart, point_chart)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Extract and process objects from PowerPoint files."
    )
    parser.add_argument(
        "-f", "--file_path", type=str, help="Path to the input file", required=True
    )
    parser.add_argument(
        "-s", "--slide_num", type=int, help="Specific slide number to process"
    )
    args = parser.parse_args()

    file_path = Path(args.file_path)
    if file_path.suffix == ".pptx":
        bin_files = extract_prism_files(file_path, single_slide=args.slide_num)

        for bin_file in bin_files:
            extracted_set = ExtractedBin(bin_file)
            extracted_set.extract_table()
            extracted_set.save_tables()
