import sys
import io
import pathlib
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import zipfile


def extract_table(filepath):
    prs = Presentation(filepath)
    input_file_directory = pathlib.Path(filepath).parent
    extracted_files = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # print page number

            # print(shape.shape_type)
            # get data from the chart
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                print(prs.slides.index(slide) + 1, "\n")
                chart = shape.chart
                print(chart.chart_title.text_frame.text)
                chart_data = []
                for series in chart.series:
                    # print attrs of series
                    print(series.name)
                    for value, category in zip(
                        series.values, chart.plots[0].categories
                    ):
                        # print(dir(point.data_label.text_frame.paragraphs))
                        # print point data label

                        print(category, value, series.name)

    return extracted_files


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: extract_obj <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    extract_table(file_path)
